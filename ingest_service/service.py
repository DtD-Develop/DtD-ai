import mimetypes
import os
from pathlib import Path
from typing import List, Optional

from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from qdrant_client import QdrantClient
from qdrant_client.http import models as qm
from sentence_transformers import SentenceTransformer

# =========================
# Config
# =========================

MODEL_NAME = "intfloat/multilingual-e5-large"
EMBED_DIM = 1024  # dim ของ multilingual-e5-large

QDRANT_HOST = os.getenv("QDRANT_HOST", "http://qdrant:6333")
QDRANT_API_KEY = os.getenv("QDRANT_API_KEY") or None
QDRANT_COLLECTION = os.getenv("QDRANT_COLLECTION", "dtd_kb")

app = FastAPI(title="DtD Ingest Service")

# โหลด model ครั้งเดียว
model = SentenceTransformer(MODEL_NAME)

# Qdrant client
qdrant = QdrantClient(
    url=QDRANT_HOST,
    api_key=QDRANT_API_KEY,
)

# =========================
# Pydantic models
# =========================


class ParseRequest(BaseModel):
    file_path: str


class Chunk(BaseModel):
    i: int
    text: str


class ParseResponse(BaseModel):
    tags: List[str]
    chunks: List[Chunk]


class EmbedRequest(BaseModel):
    file_path: str
    tags: Optional[List[str]] = None
    kb_file_id: Optional[int] = None


class EmbedResponse(BaseModel):
    chunks_count: int


# =========================
# Helpers: Collection
# =========================


def ensure_collection():
    if not qdrant.collection_exists(QDRANT_COLLECTION):
        qdrant.recreate_collection(
            collection_name=QDRANT_COLLECTION,
            vectors_config=qm.VectorParams(
                size=EMBED_DIM,
                distance=qm.Distance.COSINE,
            ),
        )


# =========================
# Helpers: Text extractors (Hybrid)
# =========================


def read_file(file_path: str) -> str:
    """
    อ่านไฟล์เป็น text แบบ basic ก่อน
    Hybrid:
      - .md, .txt, .json, code → plain text
      - .pptx → ดึงเป็นโครง markdown-like
      - .pdf → plain text (ภายหลังจูนเพิ่มได้)
      - .png/.jpg → OCR (optional)
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = path.suffix.lower()

    if ext in [".txt", ".md", ".json", ".py", ".js", ".ts", ".php", ".java", ".go"]:
        return path.read_text(encoding="utf-8", errors="ignore")

    if ext in [".pdf"]:
        # NOTE: คุณเลือกใช้ pymupdf หรือ pdfplumber ก็ได้
        try:
            import fitz  # pymupdf

            text = []
            with fitz.open(file_path) as doc:
                for page in doc:
                    text.append(page.get_text())
            return "\n".join(text)
        except Exception as e:
            raise RuntimeError(f"PDF parse error: {e}")

    if ext in [".docx"]:
        try:
            from docx import Document

            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception as e:
            raise RuntimeError(f"DOCX parse error: {e}")

    if ext in [".pptx"]:
        # Hybrid → แปลง slide เป็น markdown-ish
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            lines = []
            for i, slide in enumerate(prs.slides):
                lines.append(f"# Slide {i + 1}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines.append(shape.text.strip())
                lines.append("")  # blank line
            return "\n".join(lines)
        except Exception as e:
            raise RuntimeError(f"PPTX parse error: {e}")

    if ext in [".png", ".jpg", ".jpeg", ".webp"]:
        try:
            import pytesseract
            from PIL import Image

            img = Image.open(file_path)
            text = pytesseract.image_to_string(img)
            return text
        except Exception as e:
            raise RuntimeError(f"OCR parse error: {e}")

    # fallback ถ้าไม่รู้จัก → ลองอ่านเป็น text ปกติ
    return path.read_text(encoding="utf-8", errors="ignore")


def chunk_text(text: str, max_chars: int = 1000, overlap: int = 200) -> List[str]:
    """
    แบ่ง text เป็น chunk แบบง่าย ๆ ตามจำนวนตัวอักษร
    """
    text = text.strip()
    if not text:
        return []

    chunks = []
    start = 0
    length = len(text)

    while start < length:
        end = min(start + max_chars, length)
        chunk = text[start:end]

        # extend to nearest newline ถ้าเจอ
        newline_pos = chunk.rfind("\n")
        if newline_pos != -1 and end != length:
            end = start + newline_pos
            chunk = text[start:end]

        chunks.append(chunk.strip())
        start = max(0, end - overlap)

        if end == length:
            break

    return [c for c in chunks if c]


def auto_tag(text: str, max_tags: int = 5) -> List[str]:
    """
    Auto tag แบบเบื้องต้น:
      - split เป็นคำ
      - ตัด stopwords ง่าย ๆ
      - หาคำความถี่สูง (ความยาว > 4)
    ภายหลังคุณสามารถใช้ model ทำ zero-shot classification ได้
    """
    import re
    from collections import Counter

    text = text.lower()
    tokens = re.findall(r"[a-zA-Zก-๙]{3,}", text)

    stop = set(
        [
            "this",
            "that",
            "with",
            "from",
            "have",
            "there",
            "about",
            "your",
            "their",
            "การ",
            "และ",
            "ของ",
            "ที่",
            "ใน",
            "เป็น",
            "ได้",
            "ว่า",
        ]
    )

    words = [t for t in tokens if t not in stop and len(t) > 4]
    freq = Counter(words)
    common = [w for w, _ in freq.most_common(max_tags)]

    return common


# =========================
# Endpoints
# =========================


@app.post("/parse", response_model=ParseResponse)
def parse(req: ParseRequest):
    """
    รับ file_path จาก Laravel
    → อ่านไฟล์
    → Hybrid extract
    → chunk
    → auto-tags
    """
    try:
        text = read_file(req.file_path)
    except FileNotFoundError as e:
        raise HTTPException(status_code=404, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

    if not text.strip():
        raise HTTPException(status_code=400, detail="Empty content")

    chunks = chunk_text(text)
    tags = auto_tag(text)

    return ParseResponse(
        tags=tags,
        chunks=[Chunk(i=i, text=c) for i, c in enumerate(chunks)],
    )


@app.post("/embed", response_model=EmbedResponse)
def embed(req: EmbedRequest):
    """
    รับ file_path + tags + kb_file_id
    → extract text + chunk เหมือน parse (หรือจะ reuse /parse ก็ได้)
    → สร้าง embeddings ด้วย multilingual-e5-large
    → upsert เข้า Qdrant
    """
    try:
        text = read_file(req.file_path)
    except FileNotFoundError as e:
        raise HTTPException(status_code=404, detail=str(e))
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

    chunks = chunk_text(text)
    if not chunks:
        raise HTTPException(status_code=400, detail="No chunks generated")

    ensure_collection()

    # สร้าง embeddings
    embeddings = model.encode(chunks, batch_size=16, show_progress_bar=False)

    # เตรียม payload per chunk
    payloads = []
    for i, chunk_text_value in enumerate(chunks):
        payload = {
            "kb_file_id": req.kb_file_id,
            "chunk_index": i,
            "tags": req.tags or [],
            "text": chunk_text_value,
        }
        payloads.append(payload)

    # ใช้ upsert batch
    qdrant.upsert(
        collection_name=QDRANT_COLLECTION,
        points=qm.Batch(
            ids=[
                f"{req.kb_file_id}_{i}" if req.kb_file_id is not None else i
                for i in range(len(chunks))
            ],
            vectors=embeddings.tolist(),
            payloads=payloads,
        ),
    )

    return EmbedResponse(chunks_count=len(chunks))
